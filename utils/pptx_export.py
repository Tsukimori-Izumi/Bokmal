"""PowerPoint Export - Render Gantt chart to a single PowerPoint slide."""

from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def export_gantt_to_pptx(tasks: list[dict], dependencies: list[dict],
                          filepath: str, project_name: str = ""):
    """Export Gantt chart to a single PowerPoint slide."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Calculate date range
    all_starts = [t["start_date"] for t in tasks if isinstance(t.get("start_date"), date)]
    all_ends = [t["end_date"] for t in tasks if isinstance(t.get("end_date"), date)]
    if not all_starts or not all_ends:
        prs.save(filepath)
        return

    proj_start = min(all_starts) - timedelta(days=1)
    proj_end = max(all_ends) + timedelta(days=3)
    total_days = (proj_end - proj_start).days

    # Layout constants
    left_margin = Inches(0.3)
    top_margin = Inches(1.0)
    name_col_width = Inches(3.0)
    chart_left = left_margin + name_col_width + Inches(0.1)
    chart_width = SLIDE_WIDTH - chart_left - Inches(0.3)
    header_height = Inches(0.5)

    # 1ページに収めるためのフォールディング処理
    available_height = SLIDE_HEIGHT - top_margin - header_height - Inches(0.3)
    min_row_height = Inches(0.20)
    max_target_rows = int(available_height / min_row_height)

    display_tasks = tasks
    if len(display_tasks) > max_target_rows:
        max_level = max((t.get("wbs_level", 0) for t in display_tasks), default=0)
        for level_limit in range(max_level - 1, -1, -1):
            filtered = [t for t in display_tasks if t.get("wbs_level", 0) <= level_limit]
            if len(filtered) <= max_target_rows:
                display_tasks = filtered
                break
        else:
            display_tasks = [t for t in display_tasks if t.get("wbs_level", 0) == 0]

    if len(display_tasks) > 0:
        row_height = min(Inches(0.35), available_height / len(display_tasks))
    else:
        row_height = Inches(0.28)

    day_width = chart_width / total_days if total_days > 0 else Inches(0.1)

    # Colors
    bg_color = RGBColor(0x1A, 0x1B, 0x2E)
    header_bg = RGBColor(0x22, 0x23, 0x3C)
    bar_color = RGBColor(0x6C, 0x63, 0xFF)
    bar_progress = RGBColor(0x4E, 0xC9, 0xB0)
    critical_color = RGBColor(0xFF, 0x6B, 0x6B)
    milestone_color = RGBColor(0xFF, 0xD9, 0x3D)
    summary_color = RGBColor(0x9D, 0x8C, 0xEF)
    text_color = RGBColor(0xE0, 0xE0, 0xF0)
    grid_color = RGBColor(0x2D, 0x2E, 0x4A)

    # Slide background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Title
    txBox = slide.shapes.add_textbox(left_margin, Inches(0.2), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = project_name or "プロジェクトスケジュール"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = text_color

    # Month headers
    current_month = None
    for day_idx in range(total_days):
        d = proj_start + timedelta(days=day_idx)
        if d.month != current_month:
            current_month = d.month
            x = chart_left + day_width * day_idx
            # Find how many days until next month
            remaining = 0
            for j in range(day_idx, total_days):
                dj = proj_start + timedelta(days=j)
                if dj.month != current_month:
                    break
                remaining += 1
            w = day_width * remaining
            month_names = ["", "1月", "2月", "3月", "4月", "5月", "6月",
                           "7月", "8月", "9月", "10月", "11月", "12月"]
            label = f"{d.year}年{month_names[d.month]}"

            shape = slide.shapes.add_shape(
                1, x, top_margin, w, header_height  # MSO_SHAPE.RECTANGLE = 1
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = header_bg
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(8)
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER

    # Task rows
    base_font_size = Pt(8) if row_height >= Inches(0.25) else Pt(6)
    for row, task in enumerate(display_tasks):
        y = top_margin + header_height + row * row_height
        wbs_level = task.get("wbs_level", 0)
        indent = "    " * wbs_level

        # Task name
        name_shape = slide.shapes.add_textbox(
            left_margin, y, name_col_width, row_height
        )
        tf = name_shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        name_text = indent + task.get("name", "")
        p.text = name_text
        p.font.size = base_font_size
        if task.get("is_summary"):
            p.font.bold = True
            p.font.color.rgb = summary_color
        elif task.get("is_critical"):
            p.font.color.rgb = critical_color
        else:
            p.font.color.rgb = text_color

        # Task bar
        start = task.get("start_date")
        end = task.get("end_date")
        if not isinstance(start, date) or not isinstance(end, date):
            continue

        bar_x = chart_left + day_width * (start - proj_start).days
        bar_w = max(Emu(10000), day_width * (end - start).days)
        bar_h = row_height * 0.5
        bar_y = y + (row_height - bar_h) / 2

        if task.get("is_milestone"):
            # Diamond
            size = row_height * 0.5
            ms = slide.shapes.add_shape(
                4, bar_x - size / 2, bar_y, size, size  # MSO_SHAPE.DIAMOND = 4
            )
            ms.fill.solid()
            ms.fill.fore_color.rgb = milestone_color
            ms.line.fill.background()
        elif task.get("is_summary"):
            # Summary bar
            shape = slide.shapes.add_shape(1, bar_x, bar_y, bar_w, bar_h * 0.5)
            shape.fill.solid()
            shape.fill.fore_color.rgb = summary_color
            shape.line.fill.background()
        else:
            # Normal bar background
            shape = slide.shapes.add_shape(1, bar_x, bar_y, bar_w, bar_h)
            shape.fill.solid()
            color = critical_color if task.get("is_critical") else bar_color
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            # Progress fill
            progress = task.get("progress", 0)
            if progress and progress > 0:
                prog_w = int(bar_w * progress / 100)
                if prog_w > 0:
                    pshape = slide.shapes.add_shape(
                        1, bar_x, bar_y, prog_w, bar_h
                    )
                    pshape.fill.solid()
                    pshape.fill.fore_color.rgb = bar_progress
                    pshape.line.fill.background()

    # Today line
    today = date.today()
    if proj_start <= today <= proj_end:
        today_x = chart_left + day_width * (today - proj_start).days
        line = slide.shapes.add_shape(
            1, today_x, top_margin, Emu(15000),
            header_height + len(display_tasks) * row_height
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xFF, 0x44, 0x44)
        line.line.fill.background()

    prs.save(filepath)
